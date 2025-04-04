import os
import tempfile
from fastapi import APIRouter, UploadFile, File, HTTPException
from pydantic import BaseModel
import pdfplumber
from pptx import Presentation
from docx import Document
import librosa
import torch
from transformers import Wav2Vec2ForCTC, Wav2Vec2Processor
import requests

router = APIRouter(prefix="/api/files", tags=["files"])

def extract_text_from_pdf(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def extract_text_from_pptx(pptx_path):
    text = ""
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(docx_path):
    text = ""
    doc = Document(docx_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def transcribe_audio(audio_path):
    waveform, _ = librosa.load(audio_path, sr=16000)
    waveform = torch.tensor(waveform, dtype=torch.float32)
    model_name = "facebook/wav2vec2-large-960h"
    processor = Wav2Vec2Processor.from_pretrained(model_name)
    model = Wav2Vec2ForCTC.from_pretrained(model_name)
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    model.to(device)
    with torch.no_grad():
        inputs = processor(waveform, sampling_rate=16000, return_tensors="pt", padding=True).input_values.to(device)
        logits = model(inputs).logits
        predicted_ids = torch.argmax(logits, dim=-1)
        transcription = processor.batch_decode(predicted_ids)
    return transcription[0]

def translate_document_text_openai(text):
    """
    Utilise l'API OpenAI pour traduire le texte en français.
    Le prompt demande une traduction naturelle sans explication ni commentaire.
    """
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise Exception("OPENAI_API_KEY is not defined in the environment.")
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    data = {
        "model": "gpt-3.5-turbo",
        "messages": [
            {"role": "system", "content": "You are a translation assistant."},
            {"role": "user", "content": (
                "Translate the following English text into French. Do not translate word by word; provide a natural, fluent translation that respects proper punctuation and grammar. "
                "Return only the translated text without any additional commentary or explanation:\n\n"
                f"'''{text}'''"
            )}
        ],
        "temperature": 0.3,
        "max_tokens": 512
    }
    response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=data)
    if response.status_code == 200:
        result = response.json()
        return result["choices"][0]["message"]["content"].strip()
    else:
        raise Exception(f"Erreur API OpenAI: {response.status_code} - {response.text}")

class FileTranslationResponse(BaseModel):
    translatedText: str

@router.post("/translateFile", response_model=FileTranslationResponse)
async def translate_file(file: UploadFile = File(...)):
    suffix = os.path.splitext(file.filename)[1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        contents = await file.read()
        tmp.write(contents)
        tmp_path = tmp.name
    try:
        if suffix == ".pdf":
            text = extract_text_from_pdf(tmp_path)
        elif suffix == ".pptx":
            text = extract_text_from_pptx(tmp_path)
        elif suffix == ".docx":
            text = extract_text_from_docx(tmp_path)
        elif suffix == ".wav":
            text = transcribe_audio(tmp_path)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type.")
        translated = translate_document_text_openai(text)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        os.remove(tmp_path)
    return {"translatedText": translated}
