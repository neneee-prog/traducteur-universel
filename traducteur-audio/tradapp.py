# traducteur-audio/tradapp.py
import os
import torch
import librosa
import pdfplumber
from pptx import Presentation
from docx import Document
from transformers import Wav2Vec2ForCTC, Wav2Vec2Processor
from argostranslate import translate as argos_translate

def extract_text_from_pdf(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def extract_text_from_pptx(pptx_path):
    text = ""
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(docx_path):
    text = ""
    doc = Document(docx_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def argos_translate_text(text, source_lang="en", target_lang="fr"):
    installed_languages = argos_translate.load_installed_languages()
    from_lang = None
    to_lang = None
    for lang in installed_languages:
        if lang.code == source_lang:
            from_lang = lang
        if lang.code == target_lang:
            to_lang = lang
    if from_lang is None or to_lang is None:
        return text
    translation = from_lang.get_translation(to_lang)
    return translation.translate(text)

def transcribe_audio(audio_path):
    waveform, _ = librosa.load(audio_path, sr=16000)
    waveform = torch.tensor(waveform, dtype=torch.float32)
    processor = Wav2Vec2Processor.from_pretrained("facebook/wav2vec2-large-960h")
    model = Wav2Vec2ForCTC.from_pretrained("facebook/wav2vec2-large-960h")
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    model.to(device)
    with torch.no_grad():
        inputs = processor(waveform, sampling_rate=16000, return_tensors="pt", padding=True).input_values.to(device)
        logits = model(inputs).logits
        predicted_ids = torch.argmax(logits, dim=-1)
        transcription = processor.batch_decode(predicted_ids)
    return transcription[0]

def save_to_text_file(text, output_path):
    with open(output_path, "w", encoding="utf-8") as file:
        file.write(text)

def process_input(file_path, output_path):
    if file_path.endswith(".pdf"):
        text = extract_text_from_pdf(file_path)
        translated_text = argos_translate_text(text, "en", "fr")
        save_to_text_file(translated_text, output_path)
        return "Translation saved to text file."
    elif file_path.endswith(".pptx"):
        text = extract_text_from_pptx(file_path)
        translated_text = argos_translate_text(text, "en", "fr")
        save_to_text_file(translated_text, output_path)
        return "Translation saved to text file."
    elif file_path.endswith(".docx"):
        text = extract_text_from_docx(file_path)
        translated_text = argos_translate_text(text, "en", "fr")
        save_to_text_file(translated_text, output_path)
        return "Translation saved to text file."
    elif file_path.endswith(".wav"):
        transcription = transcribe_audio(file_path)
        save_to_text_file(transcription, output_path)
        return "Transcription saved to text file."
    else:
        return "Unsupported file type."

if __name__ == "__main__":
    result = process_input("example.pdf", "output.txt")
    print(result)